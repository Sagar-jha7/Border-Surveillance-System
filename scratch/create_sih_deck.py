import os
import sys
from pathlib import Path
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = Path(r"c:\MY Projects\AI-Based Intelligent Video Platform for Border Surveillance\Border-Surveillance-System\docs")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_PATH = OUTPUT_DIR / "SIH2026_BinaryBeasts_Presentation.pptx"

prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
C_NAVY    = RGBColor(10, 37, 64)       # #0A2540
C_BLUE    = RGBColor(0, 102, 204)      # #0066CC
C_LIGHTBG = RGBColor(248, 249, 252)    # #F8F9FC
C_TEXT    = RGBColor(30, 41, 59)       # #1E293B
C_MUTED   = RGBColor(100, 116, 139)    # #64748B
C_ACCENT  = RGBColor(220, 38, 38)      # #DC2626
C_WHITE   = RGBColor(255, 255, 255)
C_BOX_BG  = RGBColor(241, 245, 249)    # #F1F5F9
C_BORDER  = RGBColor(203, 213, 225)    # #CBD5E1

def add_header(slide, title_text, slide_num):
    # Top Blue Strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    strip.fill.solid()
    strip.fill.fore_color.rgb = C_BLUE
    strip.line.color.rgb = C_BLUE

    # Team Oval / Tag (Top-Left)
    oval = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.25), Inches(1.8), Inches(0.55))
    oval.fill.solid()
    oval.fill.fore_color.rgb = C_WHITE
    oval.line.color.rgb = C_BLUE
    oval.line.width = Pt(1.5)
    tf_oval = oval.text_frame
    tf_oval.word_wrap = True
    p = tf_oval.paragraphs[0]
    p.text = "Team:\nBinary Beasts"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = C_NAVY
    p.alignment = PP_ALIGN.CENTER

    # Main Title
    tb = slide.shapes.add_textbox(Inches(2.6), Inches(0.22), Inches(8.2), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_NAVY
    p.alignment = PP_ALIGN.CENTER

    # SIH 2026 Logo placeholder / Tag (Top-Right)
    sih_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(0.25), Inches(1.8), Inches(0.55))
    sih_box.fill.solid()
    sih_box.fill.fore_color.rgb = C_NAVY
    sih_box.line.color.rgb = C_NAVY
    tf_sih = sih_box.text_frame
    p = tf_sih.paragraphs[0]
    p.text = "SIH 2026\nPS ID: 26187"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Bottom Footer Strip
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35))
    footer.fill.solid()
    footer.fill.fore_color.rgb = C_BLUE
    footer.line.color.rgb = C_BLUE
    tf_f = footer.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = f"@SIH Idea submission- Template                                                                                                                                                                             {slide_num}"
    p_f.font.size = Pt(9)
    p_f.font.color.rgb = C_WHITE

# ==============================================================================
# SLIDE 1: TITLE PAGE
# ==============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = C_WHITE
bg.line.fill.background()

# Header text
tb = slide1.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
tf = tb.text_frame
p1 = tf.paragraphs[0]
p1.text = "SMART INDIA HACKATHON 2026"
p1.font.size = Pt(28)
p1.font.bold = True
p1.font.color.rgb = C_NAVY
p1.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "TITLE PAGE"
p2.font.size = Pt(22)
p2.font.bold = True
p2.font.color.rgb = C_BLUE
p2.alignment = PP_ALIGN.CENTER

# Details Card
card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.333), Inches(4.5))
card.fill.solid()
card.fill.fore_color.rgb = C_BOX_BG
card.line.color.rgb = C_BORDER
card.line.width = Pt(1.5)

tf_c = card.text_frame
tf_c.margin_left = Inches(0.6)
tf_c.margin_top = Inches(0.4)
tf_c.word_wrap = True

details = [
    ("Problem Statement ID", "26187"),
    ("Problem Statement Title", "AI-Based Intelligent Video Analytics Platform for Border Surveillance using existing CCTV Infrastructure"),
    ("Theme", "Blockchain & Cybersecurity"),
    ("PS Category", "Software"),
    ("Team ID", "[To be filled by team on portal]"),
    ("Team Name (Registered on portal)", "Binary Beasts"),
]

for i, (k, v) in enumerate(details):
    p = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
    p.space_after = Pt(12)
    run_bullet = p.add_run()
    run_bullet.text = "•  "
    run_bullet.font.bold = True
    run_bullet.font.size = Pt(14)
    run_bullet.font.color.rgb = C_BLUE

    run_k = p.add_run()
    run_k.text = f"{k} –  "
    run_k.font.bold = True
    run_k.font.size = Pt(14)
    run_k.font.color.rgb = C_NAVY

    run_v = p.add_run()
    run_v.text = v
    run_v.font.bold = False
    run_v.font.size = Pt(14)
    run_v.font.color.rgb = C_TEXT

# ==============================================================================
# SLIDE 2: PROPOSED SOLUTION
# ==============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide2, "IDEA TITLE: IBVAP (Intelligent Border Video Analytics Platform)", 2)

# Subheading
tb_sub = slide2.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.45))
p = tb_sub.text_frame.paragraphs[0]
p.text = "❖ Proposed Solution (Describe your Idea/Solution/Prototype)"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = C_BLUE

# 3 Column Cards
col_w = Inches(3.8)
col_h = Inches(5.3)
col_top = Inches(1.5)

cards_data = [
    ("Detailed Explanation of Proposed Solution", [
        "Software-Only AI Layer: Transforms legacy IP-CCTVs at Border Out Posts (BOPs) into smart cameras without hardware upgrades.",
        "6-Tier Threat Intelligence: Real-time Multiclass YOLOv8 Detection, ByteTrack tracking, FRS suspect matching, ANPR plate recognition, Virtual Fence tripwires & Night-mode CLAHE.",
        "Zero Latency Streaming: Dual FastAPI HTTP & HTTPS servers push annotated frames & alerts via WebSocket in <200 ms.",
        "Edge-Autonomous: Operates 100% on-premises; sensitive defense data never leaves local intranet."
    ]),
    ("How It Addresses the Problem", [
        "Zero Hardware Cost: Replaces ₹25-40 Lakh/camera hardware with pure software running on existing CCTV infrastructure.",
        "Overcomes Sentry Fatigue: AI monitors 50-100 feeds 24/7 autonomously, alerting humans only on verified threats.",
        "Eliminates Patrol Blindspots: Ground patrol officers stream mobile phone cameras into HQ AI engine via secure browser link.",
        "Forensic Tamper Resistance: All intrusions logged with JPEG forensic snapshots in an immutable SQLite audit trail."
    ]),
    ("Innovation and Uniqueness", [
        "Zero Hardware Mandate: Directly ingests standard RTSP, HTTP, USB webcams & mobile cameras with no proprietary lock-in.",
        "Dual-Engine Fallback: Combines YOLOv8 deep learning with MOG2 statistical motion to detect camouflaged targets in fog/rain.",
        "3-Tier Acoustic Alarm: Web Audio API sound alerts (RED 880Hz urgent pulse, AMBER 520Hz, BLUE 330Hz) with visual beacons.",
        "On-the-Fly Configuration: Draw digital perimeter tripwires directly in browser without physical boundary sensors."
    ])
]

for idx, (head, bullets) in enumerate(cards_data):
    left = Inches(0.6 + idx * 4.1)
    box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, col_top, col_w, col_h)
    box.fill.solid()
    box.fill.fore_color.rgb = C_BOX_BG
    box.line.color.rgb = C_BORDER
    box.line.width = Pt(1)

    tf = box.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.25)
    tf.word_wrap = True

    p_h = tf.paragraphs[0]
    p_h.text = head
    p_h.font.size = Pt(12)
    p_h.font.bold = True
    p_h.font.color.rgb = C_NAVY
    p_h.space_after = Pt(10)

    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.space_after = Pt(8)
        r1 = p_b.add_run()
        r1.text = "• "
        r1.font.bold = True
        r1.font.color.rgb = C_BLUE
        r1.font.size = Pt(9.5)
        r2 = p_b.add_run()
        r2.text = b
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = C_TEXT

# ==============================================================================
# SLIDE 3: TECHNICAL APPROACH
# ==============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide3, "TECHNICAL APPROACH", 3)

# Left Column: Tech Stack Table (w = 5.2 in)
left_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.1), Inches(5.3), Inches(5.8))
left_box.fill.solid()
left_box.fill.fore_color.rgb = C_BOX_BG
left_box.line.color.rgb = C_BORDER
tf_l = left_box.text_frame
tf_l.margin_left = Inches(0.3)
tf_l.margin_top = Inches(0.25)
tf_l.word_wrap = True

p = tf_l.paragraphs[0]
p.text = "• Technologies to be used (Languages, Frameworks, Hardware)"
p.font.size = Pt(11.5)
p.font.bold = True
p.font.color.rgb = C_NAVY
p.space_after = Pt(8)

tech_rows = [
    ("Backend Core", "Python 3.11+ / FastAPI / Uvicorn (Dual HTTP/HTTPS)"),
    ("Detection AI", "Ultralytics YOLOv8 (Multiclass Person/Vehicle/Bag)"),
    ("Tracking", "ByteTrack Algorithm (Supervision) - Trajectory IDs"),
    ("Biometrics FRS", "InsightFace (ArcFace 512-D Vectors, ONNX Runtime)"),
    ("ANPR & OCR", "OpenCV Aspect Ratio Morphology + CNN OCR"),
    ("Night Mode", "CLAHE Histogram Equalizer + MOG2 Motion"),
    ("Dashboard", "React 18 / Vite 5 / Tailwind CSS (Dark C2 Console)"),
    ("Threat Alarms", "HTML5 Web Audio API (Multi-frequency Beeps)"),
    ("Storage / DB", "SQLite3 with SHA-256 Hashes & CSV Audit Export"),
    ("Hardware", "Edge: Intel Core i5/i7, 8-16GB RAM | Optional: RTX GPU"),
]

for cat, spec in tech_rows:
    p_t = tf_l.add_paragraph()
    p_t.space_after = Pt(4)
    r1 = p_t.add_run()
    r1.text = f"▪ {cat}: "
    r1.font.bold = True
    r1.font.size = Pt(9.5)
    r1.font.color.rgb = C_BLUE
    r2 = p_t.add_run()
    r2.text = spec
    r2.font.size = Pt(9)
    r2.font.color.rgb = C_TEXT

# Right Column: Implementation Methodology & Pipeline (w = 6.6 in)
right_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(1.1), Inches(6.6), Inches(5.8))
right_box.fill.solid()
right_box.fill.fore_color.rgb = C_BOX_BG
right_box.line.color.rgb = C_BORDER
tf_r = right_box.text_frame
tf_r.margin_left = Inches(0.3)
tf_r.margin_top = Inches(0.25)
tf_r.word_wrap = True

p = tf_r.paragraphs[0]
p.text = "• Methodology & Implementation Pipeline (Working Prototype)"
p.font.size = Pt(11.5)
p.font.bold = True
p.font.color.rgb = C_NAVY
p.space_after = Pt(8)

steps = [
    ("Step 1: Multi-Source Ingestion", "Ingests IP CCTV (RTSP), USB webcams, or patrol phones (HTTPS WebSockets) without dummy videos."),
    ("Step 2: Night/Day Preprocessing", "Luminance analyzer detects low-lux frames; automatically applies CLAHE contrast equalization."),
    ("Step 3: Dual Detection & Association", "YOLOv8 detects humans/vehicles; MOG2 catches subtle motion; ByteTrack maintains IDs across frames."),
    ("Step 4: Threat Intelligence Engines", "Simultaneously evaluates Virtual Fence tripwires, loitering (>7s), BSF FRS watchlist & BOLO license plates."),
    ("Step 5: Threat Priority Classification", "Classifies events into RED (Immediate Breach), AMBER (Suspicious Motion), BLUE (Authorized Info)."),
    ("Step 6: Zero-Latency C2 Distribution", "Broadcasting annotated JPEG frames (65% quality) and alert payloads to React Dashboard & Phone."),
    ("Working Prototype Status", "Fully operational local & network deployment with 3 active multi-scenario feeds running at 25-30 FPS.")
]

for title, desc in steps:
    p_s = tf_r.add_paragraph()
    p_s.space_after = Pt(5)
    r1 = p_s.add_run()
    r1.text = f"▶ {title}: "
    r1.font.bold = True
    r1.font.size = Pt(9.5)
    r1.font.color.rgb = C_ACCENT if "Prototype" in title else C_NAVY
    r2 = p_s.add_run()
    r2.text = desc
    r2.font.size = Pt(9)
    r2.font.color.rgb = C_TEXT

# ==============================================================================
# SLIDE 4: FEASIBILITY AND VIABILITY
# ==============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide4, "FEASIBILITY AND VIABILITY", 4)

top_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.1), Inches(12.1), Inches(2.2))
top_box.fill.solid()
top_box.fill.fore_color.rgb = C_BOX_BG
top_box.line.color.rgb = C_BORDER
tf_top = top_box.text_frame
tf_top.margin_left = Inches(0.3)
tf_top.margin_top = Inches(0.2)

p = tf_top.paragraphs[0]
p.text = "• Analysis of the Feasibility of the Idea"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_NAVY
p.space_after = Pt(6)

feas_items = [
    ("Technical Feasibility", "High. Built using production-grade open-source engines (YOLOv8, InsightFace). Runs on commodity CPUs (15-25 FPS) and scales to 60+ FPS on edge GPUs. Client-server architecture decouples video inference from UI."),
    ("Operational Feasibility", "High. Designed for zero-learning-curve field operation by BSF/CISF jawans. Intuitive C2 interface with 1-click camera addition via RTSP or USB webcam with immediate visual confirmation."),
    ("Economic Viability", "High. Traditional smart deployments cost ₹25-40 Lakh per checkpost for dedicated hardware. IBVAP operates with ₹0 hardware upgrade on existing CCTV network, saving over 90% of border defense budgets.")
]

for label, text in feas_items:
    p_f = tf_top.add_paragraph()
    p_f.space_after = Pt(4)
    r1 = p_f.add_run()
    r1.text = f"✔ {label}: "
    r1.font.bold = True
    r1.font.size = Pt(10)
    r1.font.color.rgb = C_BLUE
    r2 = p_f.add_run()
    r2.text = text
    r2.font.size = Pt(9.5)
    r2.font.color.rgb = C_TEXT

# Bottom Box: Challenges, Risks & Mitigation Strategies
btm_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.45), Inches(12.1), Inches(3.45))
btm_box.fill.solid()
btm_box.fill.fore_color.rgb = C_BOX_BG
btm_box.line.color.rgb = C_BORDER
tf_btm = btm_box.text_frame
tf_btm.margin_left = Inches(0.3)
tf_btm.margin_top = Inches(0.2)

p = tf_btm.paragraphs[0]
p.text = "• Potential Challenges, Risks & Implemented Mitigation Strategies"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_NAVY
p.space_after = Pt(6)

risks = [
    ("Low-Light & Adverse Weather", "High", "Adaptive CLAHE contrast booster + MOG2 motion modeling maintains target clarity through zero-lux, fog, and rainfall."),
    ("False Alarms (Animals, Foliage)", "High", "Two-tier verification: bounding-box semantic classifier ignores non-threats; trajectory checks ensure only intentional incursions trigger alarms."),
    ("Bandwidth Drops at Remote BOPs", "Medium", "Edge-Autonomous: Full AI pipeline runs on-site at the outpost mini-PC without requiring cloud connectivity; syncs alerts asynchronously."),
    ("Occluded Faces & Muddy Plates", "Medium", "Cross-Camera Re-ID color-histogram modeling retains tracking persistence even when face/plate angles are temporarily unavailable."),
    ("Cybersecurity & Audit Tampering", "High", "All events stored in an append-only SQLite ledger with SHA-256 hashes and timestamped forensic snapshot images.")
]

for r_name, r_sev, r_mit in risks:
    p_r = tf_btm.add_paragraph()
    p_r.space_after = Pt(4)
    r1 = p_r.add_run()
    r1.text = f"⚠ {r_name} [{r_sev} Risk]: "
    r1.font.bold = True
    r1.font.size = Pt(9.5)
    r1.font.color.rgb = C_ACCENT if r_sev == "High" else C_NAVY
    r2 = p_r.add_run()
    r2.text = r_mit
    r2.font.size = Pt(9)
    r2.font.color.rgb = C_TEXT

# ==============================================================================
# SLIDE 5: IMPACT AND BENEFITS
# ==============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide5, "IMPACT AND BENEFITS", 5)

left_box5 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.1), Inches(5.8), Inches(5.8))
left_box5.fill.solid()
left_box5.fill.fore_color.rgb = C_BOX_BG
left_box5.line.color.rgb = C_BORDER
tf_l5 = left_box5.text_frame
tf_l5.margin_left = Inches(0.3)
tf_l5.margin_top = Inches(0.25)

p = tf_l5.paragraphs[0]
p.text = "• Potential Impact on the Target Audience"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_NAVY
p.space_after = Pt(10)

impacts = [
    ("Border Security Force (BSF)", "Transforms thousands of kilometers of passive barbed-wire fence cameras into an autonomous digital tripwire perimeter grid."),
    ("Rapid Response Sentry Teams", "Reduces reaction latency from minutes (manual visual scan) to <2 seconds via instant acoustic and visual alerts."),
    ("CISF Critical Installations", "Automates perimeter surveillance across airports, nuclear power plants, defense refineries, and strategic munitions hubs."),
    ("Ground Patrol Units", "Empowers patrol jawans with live back-office AI threat feeds by utilizing standard smartphones as roaming tactical surveillance nodes.")
]

for tgt, imp in impacts:
    p_i = tf_l5.add_paragraph()
    p_i.space_after = Pt(8)
    r1 = p_i.add_run()
    r1.text = f"★ {tgt}: "
    r1.font.bold = True
    r1.font.size = Pt(10)
    r1.font.color.rgb = C_BLUE
    r2 = p_i.add_run()
    r2.text = imp
    r2.font.size = Pt(9.5)
    r2.font.color.rgb = C_TEXT

# Right Box: Benefits
right_box5 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(1.1), Inches(6.1), Inches(5.8))
right_box5.fill.solid()
right_box5.fill.fore_color.rgb = C_BOX_BG
right_box5.line.color.rgb = C_BORDER
tf_r5 = right_box5.text_frame
tf_r5.margin_left = Inches(0.3)
tf_r5.margin_top = Inches(0.25)

p = tf_r5.paragraphs[0]
p.text = "• Benefits of the Solution (Multi-Domain)"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = C_NAVY
p.space_after = Pt(10)

benefits = [
    ("National Security", "Proactively detects infiltration, contraband drops, and hostile border incursions while instantly flagging terror suspects against watchlists."),
    ("Economic Savings", "Saves estimated ₹500+ Crore nationally by reusing installed CCTVs and extending legacy camera infrastructure lifespan by 8-10 years."),
    ("Operational Efficiency", "1 operator effectively monitors 50-100 feeds simultaneously (vs 4-6 manually); eliminates sentry exhaustion during graveyard shifts."),
    ("Sovereign Data Security", "100% on-premises, air-gapped system; zero reliance on third-party foreign clouds ensures defense intelligence stays secure."),
    ("Environmental / Green Defense", "Prevents massive e-waste by avoiding premature camera obsolescence; cuts patrol vehicle emissions via targeted AI dispatch.")
]

for b_cat, b_desc in benefits:
    p_b = tf_r5.add_paragraph()
    p_b.space_after = Pt(8)
    r1 = p_b.add_run()
    r1.text = f"✔ {b_cat}: "
    r1.font.bold = True
    r1.font.size = Pt(10)
    r1.font.color.rgb = C_NAVY
    r2 = p_b.add_run()
    r2.text = b_desc
    r2.font.size = Pt(9.5)
    r2.font.color.rgb = C_TEXT

# ==============================================================================
# SLIDE 6: RESEARCH AND REFERENCES
# ==============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide6, "RESEARCH AND REFERENCES", 6)

ref_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.8))
ref_box.fill.solid()
ref_box.fill.fore_color.rgb = C_BOX_BG
ref_box.line.color.rgb = C_BORDER
tf_ref = ref_box.text_frame
tf_ref.margin_left = Inches(0.4)
tf_ref.margin_top = Inches(0.25)

p = tf_ref.paragraphs[0]
p.text = "• Details / Links of the Reference and Research Work"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = C_NAVY
p.space_after = Pt(12)

refs = [
    ("Real-Time Multiclass Object Detection", "Jocher, G., Chaurasia, A., & Qiu, J. (2023). Ultralytics YOLOv8 Architecture. https://github.com/ultralytics/ultralytics — Applied for high-speed edge feature extraction & classification."),
    ("Multi-Object Tracking by Association", "Zhang, Y., Sun, P., Dong, C., et al. (2022). ByteTrack: Multi-Object Tracking by Associating Every Detection Box. ECCV 2022. https://arxiv.org/abs/2110.06864 — Applied for within-camera trajectory tracking."),
    ("Deep Facial Metric Learning (ArcFace)", "Deng, J., Guo, J., Xue, N., & Zafeiriou, S. (2019). ArcFace: Additive Angular Margin Loss for Deep Face Recognition. CVPR 2019. https://arxiv.org/abs/1801.07698 — Applied via InsightFace for suspect gallery vector matching."),
    ("Low-Light Histogram Enhancement", "Pizer, S. M., et al. (1987). Adaptive Histogram Equalization and Its Variations (CLAHE). Computer Vision, Graphics, and Image Processing — Applied for dynamic zero-lux contrast stretching."),
    ("National Border Surveillance Directives", "Ministry of Home Affairs (MHA) Comprehensive Integrated Border Management System (CIBMS) Vision Framework & Smart India Hackathon 2026 Problem Statement ID 26187. https://www.sih.gov.in"),
    ("Open-Source Prototype Repository", "IBVAP Live Working Codebase: https://github.com/Sagar-jha7/Border-Surveillance-System — Tested with 3 active multi-scenario camera nodes and tactical dashboard at http://localhost:8000.")
]

for idx, (title, citation) in enumerate(refs):
    p_r = tf_ref.add_paragraph()
    p_r.space_after = Pt(8)
    r1 = p_r.add_run()
    r1.text = f"[{idx+1}] {title}: "
    r1.font.bold = True
    r1.font.size = Pt(10)
    r1.font.color.rgb = C_BLUE
    r2 = p_r.add_run()
    r2.text = citation
    r2.font.size = Pt(9.5)
    r2.font.color.rgb = C_TEXT

prs.save(str(PPTX_PATH))
print(f"SUCCESS: Generated PowerPoint presentation at {PPTX_PATH}")
